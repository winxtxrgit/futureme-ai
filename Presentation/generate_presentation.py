#!/usr/bin/env python3
"""Generate the FutureMe AI project case-study deck.

Outputs an editable 16:9 PowerPoint presentation.  The PDF is exported from
the generated PPTX with Keynote so that the final artifact uses a real
presentation renderer rather than a second, separately maintained layout.
"""

from __future__ import annotations

import math
from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = Path(__file__).resolve().parent
ASSET_DIR = OUT_DIR / "assets"
PROJECT = ROOT / "03_WebApp" / "Pre_Present"
SCREEN_DIR = PROJECT / "assets" / "screenshots" / "app"

PPTX_PATH = OUT_DIR / "FutureMe_Project_Presentation.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

FONT = "Aptos"
FONT_THAI = "Sukhumvit Set"


def C(hex_string: str) -> RGBColor:
    value = hex_string.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


COLORS = {
    "canvas": C("#0B0B14"),
    "canvas2": C("#10101C"),
    "surface": C("#171724"),
    "surface2": C("#202033"),
    "stroke": C("#34344A"),
    "text": C("#F5F5FA"),
    "muted": C("#A8A8BE"),
    "muted2": C("#77778F"),
    "indigo": C("#6D5EF6"),
    "mint": C("#4FE3C1"),
    "mint_dark": C("#0B8B75"),
    "magenta": C("#C13BF0"),
    "coral": C("#FF6B6B"),
    "warning": C("#FFD37A"),
    "white": C("#FFFFFF"),
    "black": C("#07070D"),
}


def make_backgrounds() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    width, height = 1600, 900

    def radial_glow(base: str, glows: Sequence[tuple[int, int, int, str, int]], name: str) -> None:
        canvas = Image.new("RGB", (width, height), base)
        for cx, cy, radius, color, alpha in glows:
            layer = Image.new("RGBA", (width, height), (0, 0, 0, 0))
            draw = ImageDraw.Draw(layer)
            rgb = tuple(int(color[i : i + 2], 16) for i in (1, 3, 5))
            draw.ellipse(
                (cx - radius, cy - radius, cx + radius, cy + radius),
                fill=(*rgb, alpha),
            )
            layer = layer.filter(ImageFilter.GaussianBlur(radius // 2))
            canvas = Image.alpha_composite(canvas.convert("RGBA"), layer).convert("RGB")
        canvas.save(ASSET_DIR / name, quality=95)

    radial_glow(
        "#0B0B14",
        [
            (1370, 90, 430, "#6D5EF6", 95),
            (1500, 730, 330, "#C13BF0", 55),
            (120, 830, 360, "#4FE3C1", 45),
        ],
        "bg_cover.png",
    )
    radial_glow(
        "#0B0B14",
        [(1490, 100, 320, "#6D5EF6", 38), (40, 820, 240, "#4FE3C1", 22)],
        "bg_subtle.png",
    )

    accent = Image.new("RGB", (1600, 8), "#6D5EF6")
    px = accent.load()
    c1, c2, c3 = (109, 94, 246), (193, 59, 240), (255, 107, 107)
    for x in range(accent.width):
        t = x / (accent.width - 1)
        if t < 0.55:
            u = t / 0.55
            rgb = tuple(round(c1[i] * (1 - u) + c2[i] * u) for i in range(3))
        else:
            u = (t - 0.55) / 0.45
            rgb = tuple(round(c2[i] * (1 - u) + c3[i] * u) for i in range(3))
        for y in range(accent.height):
            px[x, y] = rgb
    accent.save(ASSET_DIR / "accent.png")


def set_solid_background(slide, color: RGBColor = COLORS["canvas"]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_background(slide, kind: str = "subtle") -> None:
    if kind == "plain":
        set_solid_background(slide)
        return
    image_name = "bg_cover.png" if kind == "cover" else "bg_subtle.png"
    slide.shapes.add_picture(
        str(ASSET_DIR / image_name),
        Inches(0),
        Inches(0),
        width=Inches(SLIDE_W),
        height=Inches(SLIDE_H),
    )


def add_shape(
    slide,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor | None = None,
    line: RGBColor | None = None,
    line_width: float = 1.0,
    radius: bool = False,
):
    actual_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else shape_type
    shape = slide.shapes.add_shape(
        actual_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: RGBColor | None = None,
    bold: bool = False,
    font: str = FONT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0,
    tracking: float | None = None,
    italic: bool = False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or COLORS["text"]
    if tracking is not None:
        # PowerPoint stores character spacing in 1/1000 pt.
        run._r.get_or_add_rPr().set("spc", str(int(tracking * 1000)))
    return box


def add_rich_line(
    slide,
    runs: Sequence[tuple[str, dict]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    for text, style in runs:
        run = p.add_run()
        run.text = text
        run.font.name = style.get("font", FONT)
        run.font.size = Pt(style.get("size", 18))
        run.font.bold = style.get("bold", False)
        run.font.italic = style.get("italic", False)
        run.font.color.rgb = style.get("color", COLORS["text"])
    return box


def add_paragraphs(
    slide,
    items: Sequence[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    color: RGBColor | None = None,
    bullet_color: RGBColor | None = None,
    spacing: float = 7,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(spacing)
        p.line_spacing = 1.05
        p.text = f"•  {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["muted"]
    return box


def add_pill(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    fill: RGBColor,
    text_color: RGBColor,
    line: RGBColor | None = None,
    size: float = 10.5,
    bold: bool = True,
):
    shape = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        0.32,
        fill=fill,
        line=line,
        line_width=0.8,
    )
    add_text(
        slide,
        text,
        x + 0.05,
        y + 0.015,
        w - 0.1,
        0.28,
        size=size,
        color=text_color,
        bold=bold,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return shape


def add_circle_label(
    slide,
    label: str,
    x: float,
    y: float,
    d: float,
    *,
    fill: RGBColor,
    text_color: RGBColor = COLORS["text"],
    size: float = 14,
    line: RGBColor | None = None,
):
    add_shape(slide, MSO_SHAPE.OVAL, x, y, d, d, fill=fill, line=line)
    add_text(
        slide,
        label,
        x + 0.08,
        y + 0.08,
        d - 0.16,
        d - 0.16,
        size=size,
        color=text_color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: RGBColor = COLORS["stroke"],
    width: float = 1.5,
    dash: bool = False,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def add_arrow_text(slide, x: float, y: float, color: RGBColor = COLORS["muted2"], size: float = 20):
    add_text(slide, "→", x, y, 0.35, 0.35, size=size, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_image_cover(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    focus_x: float = 0.5,
    focus_y: float = 0.5,
):
    with Image.open(path) as image:
        iw, ih = image.size
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    image_ratio = iw / ih
    box_ratio = w / h
    if image_ratio > box_ratio:
        visible = box_ratio / image_ratio
        crop = 1 - visible
        left = crop * max(0, min(1, focus_x))
        right = crop - left
        pic.crop_left = left
        pic.crop_right = right
    elif image_ratio < box_ratio:
        visible = image_ratio / box_ratio
        crop = 1 - visible
        top = crop * max(0, min(1, focus_y))
        bottom = crop - top
        pic.crop_top = top
        pic.crop_bottom = bottom
    return pic


def add_screenshot_frame(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    focus_y: float = 0.0,
    label: str | None = None,
):
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x - 0.05,
        y - 0.05,
        w + 0.1,
        h + 0.1,
        fill=COLORS["white"],
        line=COLORS["stroke"],
        line_width=1.0,
    )
    add_image_cover(slide, path, x, y, w, h, focus_y=focus_y)
    if label:
        add_pill(
            slide,
            label,
            x + 0.12,
            y + 0.12,
            min(w - 0.24, max(1.1, len(label) * 0.075)),
            fill=COLORS["canvas"],
            text_color=COLORS["text"],
            line=COLORS["stroke"],
            size=9.5,
        )


def add_chrome(slide, number: int, section: str, title: str, *, subtitle: str | None = None):
    add_text(
        slide,
        section.upper(),
        0.62,
        0.34,
        4.5,
        0.24,
        size=9.5,
        color=COLORS["mint"],
        bold=True,
        tracking=0.8,
    )
    add_text(slide, title, 0.62, 0.68, 12.0, 0.56, size=28, color=COLORS["text"], bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.64, 1.23, 11.9, 0.45, size=14.5, color=COLORS["muted"])
    add_text(
        slide,
        f"{number:02d}",
        12.15,
        0.36,
        0.55,
        0.25,
        size=10,
        color=COLORS["muted2"],
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    slide.shapes.add_picture(
        str(ASSET_DIR / "accent.png"),
        Inches(0.62),
        Inches(7.22),
        width=Inches(1.45),
        height=Inches(0.035),
    )
    add_text(
        slide,
        "FutureMe AI · Project case study · August 2026",
        2.18,
        7.08,
        5.2,
        0.22,
        size=8.5,
        color=COLORS["muted2"],
    )


def add_source(slide, text: str) -> None:
    add_text(
        slide,
        text,
        7.05,
        7.02,
        5.68,
        0.28,
        size=7.8,
        color=COLORS["muted2"],
        align=PP_ALIGN.RIGHT,
    )


def add_section_label(slide, text: str, x: float, y: float, w: float, color: RGBColor):
    add_text(slide, text.upper(), x, y, w, 0.22, size=9.5, color=color, bold=True, tracking=0.55)


def new_slide(prs: Presentation, number: int, section: str, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, "subtle")
    add_chrome(slide, number, section, title, subtitle=subtitle)
    return slide


def slide_01(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, "cover")
    add_text(
        slide,
        "JUMP THAILAND HACKATHON 2026 · PROJECT CASE STUDY",
        0.68,
        0.52,
        11.4,
        0.28,
        size=10,
        color=COLORS["mint"],
        bold=True,
        tracking=0.7,
    )
    add_text(slide, "FutureMe AI", 0.68, 1.18, 6.1, 0.86, size=48, color=COLORS["text"], bold=True)
    add_rich_line(
        slide,
        [
            ("From ", {"size": 23, "color": COLORS["muted"]}),
            ("“I think”", {"size": 23, "color": COLORS["warning"], "bold": True}),
            (" to ", {"size": 23, "color": COLORS["muted"]}),
            ("“I tried.”", {"size": 23, "color": COLORS["mint"], "bold": True}),
        ],
        0.72,
        2.12,
        6.2,
        0.42,
    )
    add_text(
        slide,
        "ช่วยให้นักเรียนเข้าใจตัวเอง ก่อนตัดสินใจอนาคต",
        0.72,
        2.68,
        6.2,
        0.48,
        size=18,
        color=COLORS["text"],
        bold=True,
        font=FONT_THAI,
    )
    add_text(
        slide,
        "A decision-support prototype for Thai secondary and vocational students.",
        0.72,
        3.35,
        5.85,
        0.55,
        size=16,
        color=COLORS["muted"],
    )
    add_pill(
        slide,
        "RESEARCH-INFORMED · NOT VALIDATED GUIDANCE",
        0.72,
        4.06,
        4.35,
        fill=COLORS["surface2"],
        text_color=COLORS["warning"],
        line=COLORS["warning"],
        size=9.4,
    )
    stats = [
        ("23,257", "หลักสูตรจริง", "ปวช. · ปวส. · ปริญญาตรี"),
        ("993", "สถาบัน", "ทั่วประเทศ"),
        ("923", "อาชีพ", "โปรไฟล์ความสนใจที่วัดมา"),
    ]
    for i, (big, label, sub_label) in enumerate(stats):
        x = 0.72 + i * 2.02
        add_text(slide, big, x, 4.78, 1.9, 0.52, size=26, color=COLORS["mint"], bold=True)
        add_text(slide, label, x, 5.34, 1.9, 0.26, size=12, color=COLORS["text"], bold=True, font=FONT_THAI)
        add_text(slide, sub_label, x, 5.62, 1.94, 0.24, size=9, color=COLORS["muted2"], font=FONT_THAI)
    add_line(slide, 0.74, 4.62, 6.5, 4.62, color=COLORS["stroke"], width=1.0)

    add_text(
        slide,
        "Explore the next step—not one final answer.",
        0.72,
        6.2,
        5.6,
        0.34,
        size=15.5,
        color=COLORS["text"],
        bold=True,
    )
    add_text(
        slide,
        "Research synthesis · Product strategy · UX · AI architecture · Validation",
        0.72,
        6.63,
        6.3,
        0.26,
        size=9.5,
        color=COLORS["muted2"],
    )

    # Original route-branch illustration.
    cx, cy = 9.55, 3.62
    add_circle_label(slide, "ME\nNOW", cx - 0.52, cy - 0.52, 1.04, fill=COLORS["white"], text_color=COLORS["black"], size=13)
    endpoints = [
        (8.55, 1.08, "STUDY\nROUTE", COLORS["indigo"]),
        (10.65, 1.14, "TVET\nROUTE", COLORS["magenta"]),
        (11.25, 4.75, "SKILL\nTRIAL", COLORS["coral"]),
    ]
    for ex, ey, label, color in endpoints:
        add_line(slide, cx, cy, ex + 0.52, ey + 0.52, color=COLORS["muted2"], width=2.0)
        add_circle_label(slide, label, ex, ey, 1.04, fill=color, text_color=COLORS["white"], size=9)
    add_line(slide, cx + 0.45, cy + 0.45, 9.0, 5.75, color=COLORS["mint"], width=2.3)
    add_circle_label(slide, "NEXT\nEVIDENCE", 8.48, 5.62, 1.15, fill=COLORS["mint"], text_color=COLORS["black"], size=11)
    add_text(
        slide,
        "Several hypotheses. One reversible next step.",
        7.45,
        6.98,
        4.9,
        0.3,
        size=12,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
    )


def slide_02(prs: Presentation):
    slide = new_slide(
        prs,
        2,
        "Problem",
        "Students decide before they have evidence",
        "Thai education has visible forks; the day-to-day reality behind each route is harder to test.",
    )
    add_text(
        slide,
        "Decisions arrive\nbefore experience.",
        0.74,
        1.82,
        4.15,
        1.0,
        size=25,
        color=COLORS["text"],
        bold=True,
    )
    add_text(
        slide,
        "FutureMe’s design hypothesis",
        0.76,
        3.0,
        3.6,
        0.26,
        size=10,
        color=COLORS["warning"],
        bold=True,
        tracking=0.45,
    )
    add_text(
        slide,
        "Students need a safe way to test a direction before committing time, prerequisites and money.",
        0.76,
        3.34,
        3.95,
        1.02,
        size=18,
        color=COLORS["muted"],
    )

    # Decision fork.
    add_section_label(slide, "CONSEQUENTIAL CHOICE POINTS", 5.35, 1.76, 4.5, COLORS["mint"])
    add_circle_label(slide, "ม.3", 5.55, 2.48, 0.82, fill=COLORS["indigo"], size=15)
    add_circle_label(slide, "ม.6", 5.55, 4.58, 0.82, fill=COLORS["magenta"], size=15)
    routes_top = [
        (7.15, 2.12, "General\nupper secondary", COLORS["surface2"]),
        (9.55, 2.12, "Vocational\nปวช.", COLORS["surface2"]),
        (11.15, 2.12, "Other /\nlocal options", COLORS["surface2"]),
    ]
    routes_bottom = [
        (7.15, 4.2, "TCAS /\nuniversity", COLORS["surface2"]),
        (9.55, 4.2, "Work /\ntraining", COLORS["surface2"]),
        (11.15, 4.2, "Continue\nvocational", COLORS["surface2"]),
    ]
    for x, y, _, _ in routes_top:
        add_line(slide, 6.36, 2.89, x, y + 0.42, color=COLORS["stroke"], width=1.5)
    for x, y, label, fill in routes_top:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 1.55, 0.82, fill=fill, line=COLORS["stroke"])
        add_text(slide, label, x + 0.08, y + 0.1, 1.39, 0.64, size=12, color=COLORS["text"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    for x, y, _, _ in routes_bottom:
        add_line(slide, 6.36, 4.99, x, y + 0.42, color=COLORS["stroke"], width=1.5)
    for x, y, label, fill in routes_bottom:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 1.55, 0.82, fill=fill, line=COLORS["stroke"])
        add_text(slide, label, x + 0.08, y + 0.1, 1.39, 0.64, size=12, color=COLORS["text"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.35, 5.62, 7.15, 0.74, fill=COLORS["surface"], line=COLORS["warning"])
    add_text(
        slide,
        "Criteria vary by school, campus, programme and admission year.",
        5.62,
        5.86,
        6.64,
        0.28,
        size=14,
        color=COLORS["warning"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.75, 5.9, 3.95, 0.6, fill=COLORS["surface2"], line=COLORS["stroke"])
    add_text(
        slide,
        "PRIMARY USER RESEARCH: NOT YET RUN",
        0.88,
        6.08,
        3.68,
        0.24,
        size=10.5,
        color=COLORS["warning"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(slide, "Sources: Thai curriculum, VEC catalogue, TCAS70; FutureMe source audit (24 Jul 2026)")


def slide_03(prs: Presentation):
    slide = new_slide(
        prs,
        3,
        "Evidence",
        "The cost is mismatch—and a moving target",
        "The evidence supports earlier exploration, not a promise of one perfect match.",
    )
    metrics = [
        ("56%", "work outside their field of study", "TDRI context", COLORS["indigo"]),
        ("27%", "work below their skill or qualification level", "same population", COLORS["magenta"]),
        ("39%", "of core skills expected to change by 2030", "WEF employer survey", COLORS["mint"]),
    ]
    xs = [0.72, 4.62, 8.52]
    for x, (value, label, scope, color) in zip(xs, metrics):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.85, 3.55, 2.65, fill=COLORS["surface"], line=COLORS["stroke"])
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, 1.85, 0.08, 2.65, fill=color)
        add_text(slide, value, x + 0.28, 2.06, 3.0, 0.82, size=44, color=color, bold=True)
        add_text(slide, label, x + 0.3, 2.93, 2.92, 0.78, size=16, color=COLORS["text"], bold=True)
        add_text(slide, scope.upper(), x + 0.3, 3.94, 2.92, 0.22, size=9.3, color=COLORS["muted2"], bold=True, tracking=0.35)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.2, 5.08, 10.92, 1.0, fill=COLORS["surface2"], line=COLORS["mint"])
    add_rich_line(
        slide,
        [
            ("Design implication: ", {"size": 18, "color": COLORS["mint"], "bold": True}),
            ("help students test, compare and revise—before a direction becomes expensive to change.", {"size": 18, "color": COLORS["text"], "bold": True}),
        ],
        1.5,
        5.4,
        10.3,
        0.52,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Caution: the TDRI public page describes a broadly highly educated group and does not expose its denominator or method. These figures are problem context, not a FutureMe performance baseline.",
        1.15,
        6.35,
        11.05,
        0.44,
        size=10.5,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
    )
    add_source(slide, "Sources: TDRI (2025); WEF Future of Jobs (2025); scopes differ and are not combined causally")


def slide_04(prs: Presentation):
    slide = new_slide(
        prs,
        4,
        "Current gap",
        "Information exists. Interpretation is fragmented.",
        "The gap is not another list—it is a visible chain from self-knowledge to action.",
    )
    y = 2.08
    stages = [
        ("SEARCH", "School plans\nTCAS pages\nCareer sites", COLORS["indigo"]),
        ("TEST", "Static interest\nquestionnaire", COLORS["magenta"]),
        ("RESULT", "Generic career\nor faculty list", COLORS["coral"]),
        ("AFTER", "Still unsure\nwhat to try", COLORS["warning"]),
    ]
    x_positions = [0.75, 3.55, 6.35, 9.15]
    for i, ((cap, body, color), x) in enumerate(zip(stages, x_positions)):
        add_text(slide, cap, x, y - 0.38, 2.2, 0.2, size=9.5, color=color, bold=True, tracking=0.45, align=PP_ALIGN.CENTER)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 2.25, 1.42, fill=COLORS["surface"], line=color, line_width=1.2)
        add_text(slide, body, x + 0.13, y + 0.22, 1.99, 0.98, size=15, color=COLORS["text"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(stages) - 1:
            add_arrow_text(slide, x + 2.37, y + 0.5, COLORS["muted2"], 22)
    add_text(
        slide,
        "What is missing",
        0.78,
        4.02,
        2.1,
        0.26,
        size=10,
        color=COLORS["mint"],
        bold=True,
        tracking=0.5,
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.75, 4.42, 11.82, 1.24, fill=COLORS["surface2"], line=COLORS["mint"], line_width=1.25)
    chain = [
        ("WHAT I SAID", COLORS["indigo"]),
        ("WHAT I TRIED", COLORS["magenta"]),
        ("WHY THIS ROUTE", COLORS["mint"]),
        ("WHAT TO DO NEXT", COLORS["warning"]),
    ]
    cx = [1.1, 4.05, 7.0, 9.95]
    for i, ((label, color), x) in enumerate(zip(chain, cx)):
        add_pill(slide, label, x, 4.82, 2.0, fill=COLORS["canvas"], text_color=color, line=color, size=9.5)
        if i < len(chain) - 1:
            add_arrow_text(slide, x + 2.15, 4.77, COLORS["muted2"], 18)
    add_text(
        slide,
        "Admission criteria expire. Career-to-programme links are many-to-many. A trustworthy product must show scope, date and source.",
        1.05,
        6.08,
        11.2,
        0.48,
        size=13.5,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
    )
    add_source(slide, "Sources: FutureMe Data/RAG_DATA_SPEC; curriculum and source-audit files")


def slide_05(prs: Presentation):
    slide = new_slide(
        prs,
        5,
        "Idea",
        "The idea evolved from matching to learning",
        "Each research input changed a product decision.",
    )
    entries = [
        ("1", "PROBLEM", "Study–work mismatch\nis not one-to-one", COLORS["coral"]),
        ("2", "THEORY", "RIASEC gives a\nsix-part interest language", COLORS["indigo"]),
        ("3", "CONTEXT", "Socratic + STAR\nsurface lived evidence", COLORS["magenta"]),
        ("4", "ACTION", "A mini mission\ntests self-report", COLORS["mint"]),
        ("5", "PRODUCT", "Top 5 programmes +\na 30-day experiment", COLORS["warning"]),
    ]
    x_positions = [0.7, 3.15, 5.6, 8.05, 10.5]
    for i, ((num, cap, body, color), x) in enumerate(zip(entries, x_positions)):
        add_circle_label(slide, num, x + 0.62, 2.0, 0.72, fill=color, text_color=COLORS["black"] if color in (COLORS["mint"], COLORS["warning"]) else COLORS["white"], size=15)
        add_text(slide, cap, x, 2.9, 1.98, 0.22, size=9.5, color=color, bold=True, tracking=0.45, align=PP_ALIGN.CENTER)
        add_text(slide, body, x, 3.28, 1.98, 0.82, size=14, color=COLORS["text"], bold=True, align=PP_ALIGN.CENTER)
        if i < len(entries) - 1:
            add_arrow_text(slide, x + 2.05, 2.15, COLORS["muted2"], 20)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.02, 5.1, 9.3, 0.92, fill=COLORS["surface2"], line=COLORS["mint"])
    add_rich_line(
        slide,
        [
            ("FutureMe is an ", {"size": 21, "color": COLORS["text"]}),
            ("evidence loop", {"size": 21, "color": COLORS["mint"], "bold": True}),
            (", not a career oracle.", {"size": 21, "color": COLORS["text"], "bold": True}),
        ],
        2.28,
        5.38,
        8.8,
        0.42,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Conversation frameworks guide prompts; they do not validate the questionnaire.",
        2.65,
        6.3,
        8.05,
        0.32,
        size=12,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
    )
    add_source(slide, "Sources: FutureMe research summaries; Holland RIASEC; interview framework files")


def slide_06(prs: Presentation):
    slide = new_slide(
        prs,
        6,
        "Research foundation",
        "Research shaped the model—and what we rejected",
        "Only constructs that change the product’s output earn a place in the assessment.",
    )
    # RIASEC hexagon.
    hex_shape = add_shape(slide, MSO_SHAPE.HEXAGON, 0.92, 2.08, 3.5, 3.5, fill=COLORS["surface"], line=COLORS["indigo"], line_width=1.6)
    add_text(slide, "RIASEC", 1.87, 3.16, 1.6, 0.5, size=24, color=COLORS["text"], bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "one vector\nsix interests", 1.97, 3.72, 1.4, 0.55, size=12, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    labels = [
        ("R", "Hands-on", 2.22, 1.67, COLORS["coral"]),
        ("I", "Investigate", 3.75, 2.55, COLORS["indigo"]),
        ("A", "Create", 3.7, 4.42, COLORS["magenta"]),
        ("S", "Help", 2.18, 5.28, COLORS["mint"]),
        ("E", "Lead", 0.52, 4.4, COLORS["warning"]),
        ("C", "Organise", 0.52, 2.53, COLORS["muted"]),
    ]
    for code, label, x, y, color in labels:
        add_circle_label(slide, code, x, y, 0.46, fill=color, text_color=COLORS["black"] if color in (COLORS["mint"], COLORS["warning"], COLORS["muted"]) else COLORS["white"], size=11)
        add_text(slide, label, x - 0.25, y + 0.5, 0.96, 0.2, size=9, color=COLORS["muted"], align=PP_ALIGN.CENTER)

    rows = [
        ("INCLUDED", "RIASEC interest vector", "A starting hypothesis for exploration", COLORS["mint"]),
        ("DESIGN INPUT", "Socratic prompts + STAR", "Conversation scaffolds—not psychometrics", COLORS["indigo"]),
        ("DEFERRED", "Work values", "Route catalogue lacks value profiles", COLORS["warning"]),
        ("REJECTED", "MBTI-type matching · learning styles", "Unstable types / weak evidence", COLORS["coral"]),
    ]
    y = 1.94
    for cap, title, reason, color in rows:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.2, y, 7.0, 0.92, fill=COLORS["surface"], line=COLORS["stroke"])
        add_text(slide, cap, 5.48, y + 0.18, 1.22, 0.2, size=9.2, color=color, bold=True, tracking=0.35)
        add_text(slide, title, 6.82, y + 0.14, 2.45, 0.28, size=15, color=COLORS["text"], bold=True)
        add_text(slide, reason, 9.28, y + 0.16, 2.62, 0.48, size=12, color=COLORS["muted"])
        y += 1.08
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.2, 6.28, 7.0, 0.48, fill=COLORS["surface2"], line=COLORS["warning"])
    add_text(
        slide,
        "The 41-item instrument is research-informed—but has never been validated.",
        5.45,
        6.42,
        6.5,
        0.2,
        size=11.5,
        color=COLORS["warning"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(slide, "Sources: Holland (1997); O*NET Interest Profiler Manual; FutureMe methodology")


def slide_07(prs: Presentation):
    slide = new_slide(
        prs,
        7,
        "Design principles",
        "Five product rules follow directly from the evidence",
        "Research is useful only when it changes what the product does.",
    )
    mappings = [
        ("Interests still move in adolescence", "Offer hypotheses—not a verdict; keep routes reversible", COLORS["indigo"]),
        ("Self-report is only one signal", "Add a scenario mission that can agree or disagree", COLORS["magenta"]),
        ("AI recommendations can become opaque", "Rules decide; show evidence, unknowns and contradictions", COLORS["mint"]),
        ("Admissions and labour data expire", "Carry source, scope, last-checked date and freshness", COLORS["warning"]),
        ("The users are minors", "Guest-first, browser-local; consent before saving or sharing", COLORS["coral"]),
    ]
    y = 1.72
    for idx, (insight, decision, color) in enumerate(mappings, start=1):
        add_text(slide, f"{idx:02d}", 0.75, y + 0.15, 0.55, 0.25, size=12, color=color, bold=True)
        add_text(slide, insight, 1.48, y + 0.08, 4.05, 0.42, size=15.5, color=COLORS["text"], bold=True)
        add_arrow_text(slide, 5.62, y + 0.03, color, 18)
        add_text(slide, decision, 6.18, y + 0.08, 5.92, 0.42, size=15, color=COLORS["muted"])
        add_line(slide, 0.75, y + 0.64, 12.1, y + 0.64, color=COLORS["stroke"], width=0.9)
        y += 0.92
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.48, 6.44, 10.1, 0.46, fill=COLORS["surface2"], line=COLORS["mint"])
    add_text(
        slide,
        "Private by default · Evidence before confidence · One useful next action",
        1.72,
        6.56,
        9.6,
        0.22,
        size=12,
        color=COLORS["mint"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(slide, "Sources: Source audit, RAG data spec, interview research, privacy design")


def slide_08(prs: Presentation):
    slide = new_slide(
        prs,
        8,
        "Solution",
        "FutureMe turns reflection into a low-stakes experiment",
        "The experience accumulates evidence instead of manufacturing certainty.",
    )
    steps = [
        ("REFLECT", "41-item profile:\ninterest + confidence", COLORS["indigo"]),
        ("TRY", "One scenario\nmission", COLORS["magenta"]),
        ("EXPLORE", "Top 5 real\nprogrammes", COLORS["mint"]),
        ("COMPARE", "Evidence · trade-offs\nunknowns · sources", COLORS["warning"]),
        ("ACT", "A reversible\n30-day plan", COLORS["coral"]),
    ]
    xs = [0.68, 3.18, 5.68, 8.18, 10.68]
    for i, ((label, body, color), x) in enumerate(zip(steps, xs)):
        add_circle_label(slide, str(i + 1), x + 0.68, 2.0, 0.66, fill=color, text_color=COLORS["black"] if color in (COLORS["mint"], COLORS["warning"]) else COLORS["white"], size=14)
        add_text(slide, label, x, 2.92, 2.0, 0.22, size=10, color=color, bold=True, tracking=0.48, align=PP_ALIGN.CENTER)
        add_text(slide, body, x, 3.32, 2.0, 0.86, size=14, color=COLORS["text"], bold=True, align=PP_ALIGN.CENTER)
        if i < 4:
            add_arrow_text(slide, x + 2.08, 2.18, COLORS["muted2"], 20)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.36, 5.0, 8.6, 1.0, fill=COLORS["surface2"], line=COLORS["mint"])
    add_text(slide, "EVIDENCE — NOT CERTAINTY", 2.75, 5.25, 7.8, 0.28, size=12, color=COLORS["mint"], bold=True, tracking=0.7, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "A learner can revise answers, override the mission, compare routes or stop.",
        2.75,
        5.62,
        7.8,
        0.42,
        size=13.2,
        color=COLORS["text"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(slide, "Source: implemented FutureMe guest journey")


def slide_09(prs: Presentation):
    slide = new_slide(
        prs,
        9,
        "User journey",
        "Each step gives the learner control",
        "The output at every stage is provisional, inspectable and editable.",
    )
    questions = [
        ("1", "What activities\ndraw me?", "Provisional\nRIASEC vector", COLORS["indigo"]),
        ("2", "Does action\nagree?", "Mission can\ncorroborate or contradict", COLORS["magenta"]),
        ("3", "What is\nfeasible now?", "Cost · location ·\ntiming filters", COLORS["warning"]),
        ("4", "Why these\noptions?", "Reasons · evidence ·\nunknowns · sources", COLORS["mint"]),
        ("5", "What can I\ntest next?", "30-day reversible\nexperiment", COLORS["coral"]),
    ]
    x_positions = [0.7, 3.18, 5.66, 8.14, 10.62]
    add_line(slide, 1.25, 2.63, 11.95, 2.63, color=COLORS["stroke"], width=2.0)
    for num, question, output, color in questions:
        x = x_positions[int(num) - 1]
        add_circle_label(slide, num, x + 0.68, 2.25, 0.76, fill=color, text_color=COLORS["black"] if color in (COLORS["mint"], COLORS["warning"]) else COLORS["white"], size=15)
        add_text(slide, question, x, 3.28, 2.12, 0.72, size=15.5, color=COLORS["text"], bold=True, align=PP_ALIGN.CENTER)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 4.32, 2.12, 0.82, fill=COLORS["surface"], line=COLORS["stroke"])
        add_text(slide, output, x + 0.12, 4.5, 1.88, 0.48, size=11.5, color=COLORS["muted"], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.18, 5.75, 11.0, 0.66, fill=COLORS["surface2"], line=COLORS["mint"])
    add_text(
        slide,
        "CONTROL LOOP  ·  change answers  ·  choose another mission  ·  compare again  ·  keep, park or drop a route",
        1.4,
        5.96,
        10.55,
        0.22,
        size=11,
        color=COLORS["mint"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(slide, "Source: implemented UX flow and session behavior")


def slide_10(prs: Presentation):
    slide = new_slide(
        prs,
        10,
        "AI architecture",
        "AI supports the conversation—not the decision",
        "Current capabilities and future architecture are deliberately separated.",
    )
    # Left: implemented.
    add_pill(slide, "TODAY · IMPLEMENTED", 0.76, 1.67, 2.05, fill=COLORS["mint"], text_color=COLORS["black"], size=9.5)
    left_nodes = [
        ("Learner input", COLORS["surface2"]),
        ("Deterministic\nTypeScript rules", COLORS["indigo"]),
        ("Top 5 programmes", COLORS["mint"]),
    ]
    ly = [2.24, 3.42, 4.6]
    for i, ((label, color), y) in enumerate(zip(left_nodes, ly)):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, y, 3.7, 0.72, fill=color, line=COLORS["stroke"])
        add_text(slide, label, 1.05, y + 0.15, 3.4, 0.42, size=14.5, color=COLORS["black"] if color == COLORS["mint"] else COLORS["text"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < 2:
            add_text(slide, "↓", 2.48, y + 0.76, 0.52, 0.32, size=20, color=COLORS["muted2"], bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.2, 5.72, 3.08, 0.72, fill=COLORS["surface"], line=COLORS["magenta"])
    add_text(slide, "Optional LLM: reword reasons only\nOff by default · no learner answers", 1.34, 5.88, 2.8, 0.4, size=11, color=COLORS["muted"], align=PP_ALIGN.CENTER)

    # Divider.
    add_line(slide, 5.03, 1.68, 5.03, 6.55, color=COLORS["stroke"], width=1.3, dash=True)

    # Right: planned.
    add_pill(slide, "FUTURE · PLANNED", 5.52, 1.67, 1.86, fill=COLORS["surface2"], text_color=COLORS["warning"], line=COLORS["warning"], size=9.5)
    planned = [
        ("Adaptive Thai conversation", COLORS["surface"]),
        ("Structured profile extraction", COLORS["surface"]),
        ("Rule engine", COLORS["indigo"]),
        ("Audited RAG retrieval", COLORS["surface"]),
        ("Evidence-grounded explanation", COLORS["surface"]),
        ("Interactive roadmap", COLORS["surface"]),
    ]
    px = [5.5, 7.55, 9.6, 5.5, 7.55, 9.6]
    py = [2.35, 2.35, 2.35, 4.15, 4.15, 4.15]
    for i, ((label, color), x, y) in enumerate(zip(planned, px, py)):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 1.8, 0.92, fill=color, line=COLORS["warning"] if color == COLORS["surface"] else COLORS["stroke"])
        add_text(slide, label, x + 0.12, y + 0.18, 1.56, 0.54, size=11.5, color=COLORS["text"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i in (0, 1, 3, 4):
            add_arrow_text(slide, x + 1.82, y + 0.27, COLORS["muted2"], 16)
    add_text(slide, "↓", 10.22, 3.32, 0.52, 0.42, size=20, color=COLORS["muted2"], bold=True, align=PP_ALIGN.CENTER)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.52, 5.75, 6.03, 0.68, fill=COLORS["surface2"], line=COLORS["mint"])
    add_rich_line(
        slide,
        [
            ("Rules decide. ", {"size": 16, "color": COLORS["mint"], "bold": True}),
            ("Retrieval grounds. ", {"size": 16, "color": COLORS["warning"], "bold": True}),
            ("AI explains.", {"size": 16, "color": COLORS["text"], "bold": True}),
        ],
        5.8,
        5.98,
        5.5,
        0.3,
        align=PP_ALIGN.CENTER,
    )
    add_source(slide, "Sources: current code + planned RAG architecture; Qdrant/AIS/NDLP are not implemented")


def slide_11(prs: Presentation):
    slide = new_slide(
        prs,
        11,
        "Recommendation logic",
        "Every number is measured, derived, or declared missing",
        "Same answers, same programmes. No model touches the ranking.",
    )
    pipeline = [
        ("41 items\n+ context", COLORS["surface2"]),
        ("RIASEC\nvector", COLORS["indigo"]),
        ("Kelley\nshrinkage", COLORS["magenta"]),
        ("Core fit\ncos + efficacy", COLORS["mint"]),
        ("Context\ncapped +15", COLORS["warning"]),
        ("Top 5\nprogrammes", COLORS["coral"]),
    ]
    xs = [0.66, 2.78, 4.9, 7.02, 9.14, 11.26]
    for i, ((label, color), x) in enumerate(zip(pipeline, xs)):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.82, 1.45, 0.92, fill=color, line=COLORS["stroke"])
        add_text(slide, label, x + 0.08, 2.02, 1.29, 0.5, size=11.5,
                 color=COLORS["black"] if color in (COLORS["mint"], COLORS["warning"]) else COLORS["text"],
                 bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < 5:
            add_arrow_text(slide, x + 1.52, 2.1, COLORS["muted2"], 16)

    add_section_label(slide, "INTEREST VECTORS ARE MEASURED, NOT ASSIGNED", 0.75, 3.17, 9.4, COLORS["mint"])
    add_text(
        slide,
        "A programme’s six-dimension profile is the mean of the O*NET occupations in its official ISCED-F field — "
        "923 occupations, measured. The twelve hand-written vectors it replaced had the wrong dominant dimension "
        "in four of nine comparable fields.",
        0.76, 3.56, 11.8, 0.66, size=12.5, color=COLORS["muted"],
    )

    add_section_label(slide, "ACADEMIC FIT AND CONTEXT STAY APART", 0.75, 4.34, 8.2, COLORS["warning"])
    add_text(
        slide,
        "Final = CoreFit + 15 × ContextFit.  Distance, cost and intake reorder programmes of similar fit; "
        "they cannot lift a poor fit above a strong one, because 15 points is smaller than any gap that matters.",
        0.76, 4.72, 11.8, 0.5, size=12.5, color=COLORS["muted"],
    )

    guardrails = [
        ("“3” = 0", "an unsure answer is not evidence"),
        ("FLAT", "profile → no ranking (Holland)"),
        ("NULL", "tuition unknown → link, not a guess"),
        ("TRACE", "every figure opens to its source"),
    ]
    gx = [0.75, 3.8, 6.85, 9.9]
    for (cap, body), x in zip(guardrails, gx):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 5.42, 2.68, 1.06, fill=COLORS["surface"], line=COLORS["stroke"])
        add_text(slide, cap, x + 0.16, 5.60, 2.36, 0.28, size=15, color=COLORS["mint"], bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.16, 5.96, 2.36, 0.34, size=10, color=COLORS["muted"], align=PP_ALIGN.CENTER)

    add_source(slide, "Source: lib/recommend · the Python reference in 01_Research/Recommendation_Engine is pinned to it by test")


def slide_11b(prs: Presentation):
    slide = new_slide(
        prs,
        12,
        "Data",
        "23,257 real programmes, and the gaps named out loud",
        "Every figure traces to a government register. Nothing on a card is an estimate.",
    )
    tiles = [
        ("23,257", "หลักสูตรจริง", "ปวช. 16,908 · ปวส. · ป.ตรี 6,349", COLORS["mint"]),
        ("993", "สถาบัน", "840 วิทยาลัยอาชีวะ · 153 มหาวิทยาลัย", COLORS["indigo"]),
        ("923", "อาชีพ O*NET", "RIASEC วัดจริง · 269 สาย", COLORS["magenta"]),
        ("15,586", "มีผลการมีงานทำ", "สอศ. 2566 · แสดงฐานเสมอ", COLORS["coral"]),
    ]
    xs = [0.75, 3.94, 7.13, 10.32]
    for (big, label, sub, color), x in zip(tiles, xs):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.78, 2.86, 1.62, fill=COLORS["surface"], line=COLORS["stroke"])
        add_text(slide, big, x + 0.18, 1.96, 2.5, 0.5, size=27, color=color, bold=True)
        add_text(slide, label, x + 0.18, 2.52, 2.5, 0.28, size=12.5, color=COLORS["text"], bold=True, font=FONT_THAI)
        add_text(slide, sub, x + 0.18, 2.86, 2.5, 0.42, size=10, color=COLORS["muted"], font=FONT_THAI)

    add_section_label(slide, "SOURCES · ALL GOVERNMENT REGISTERS", 0.75, 3.68, 8.0, COLORS["mint"])
    sources = [
        "แผนการรับนักศึกษา (data.go.th) · หลักสูตรและรหัส ISCED-F (อว.) · ต้นทุนต่อหัวรายหลักสูตร (อว.)",
        "นักเรียนอาชีวะรายวิทยาลัยรายสาขา (สอศ. 2568) · ภาวะการมีงานทำ (สอศ. 2566) · ระยะทางถนนจริง OSRM 77 จังหวัด",
        "O*NET 29.1 Interests (US DOL, CC BY 4.0) · ทะเบียนที่ตั้งสถานศึกษา อว. และ สอศ.",
    ]
    for i, line in enumerate(sources):
        add_text(slide, line, 0.76, 4.06 + i * 0.34, 11.8, 0.3, size=11.5, color=COLORS["muted"], font=FONT_THAI)

    add_section_label(slide, "WHAT WE DO NOT HAVE — AND DO NOT ESTIMATE", 0.75, 5.22, 9.6, COLORS["warning"])
    gaps = [
        ("ค่าเทอมจริง", "ไม่มีชุดข้อมูลเปิด → ลิงก์ประกาศของสถาบัน"),
        ("รอบ TCAS · คะแนน", "mytcas ไม่มี API → ยังตอบไม่ได้ บอกตรง ๆ"),
        ("ทุนการศึกษา", "ไม่มีทะเบียนกลาง"),
    ]
    gx = [0.75, 4.75, 8.75]
    for (cap, body), x in zip(gaps, gx):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 5.62, 3.7, 0.94, fill=COLORS["surface"], line=COLORS["warning"])
        add_text(slide, cap, x + 0.18, 5.76, 3.3, 0.28, size=13, color=COLORS["warning"], bold=True, font=FONT_THAI)
        add_text(slide, body, x + 0.18, 6.08, 3.34, 0.38, size=10, color=COLORS["muted"], font=FONT_THAI)

    add_source(slide, "A blank field is a finding. Filling it with an average would be the failure this project keeps arguing against.")


def slide_12(prs: Presentation):
    slide = new_slide(
        prs,
        13,
        "Prototype",
        "A complete prototype already runs end to end",
        "These are screenshots from the implemented app—not concept mockups.",
    )
    frames = [
        (SCREEN_DIR / "interview-desktop.png", 0.72, "1 · ASSESS", "One item at a time"),
        (SCREEN_DIR / "routes-desktop.png", 4.57, "3 · EXPLORE", "Several hypotheses"),
        (SCREEN_DIR / "plan-desktop.png", 8.42, "5 · ACT", "30-day experiment"),
    ]
    for path, x, label, caption in frames:
        add_screenshot_frame(slide, path, x, 1.77, 3.52, 3.56, focus_y=0.0, label=label)
        add_text(slide, caption, x, 5.47, 3.52, 0.24, size=12.5, color=COLORS["text"], bold=True, align=PP_ALIGN.CENTER)
    pills = [
        ("TH / EN + THEMES", 0.78, 1.85),
        ("3 MISSIONS", 2.78, 1.4),
        ("23,257 PROGRAMMES", 4.38, 2.15),
        ("BROWSER-LOCAL", 6.93, 1.65),
        ("559 TESTS", 8.58, 1.15),
        ("993 INSTITUTIONS", 9.93, 1.95),
    ]
    x = 0.78
    for text, _, width in pills:
        add_pill(slide, text, x, 6.15, width, fill=COLORS["surface2"], text_color=COLORS["mint"], line=COLORS["stroke"], size=8.8)
        x += width + 0.18
    add_text(
        slide,
        "Programmes and distances come from government registers. The interest instrument itself is still unvalidated.",
        3.2,
        6.64,
        6.95,
        0.26,
        size=10.5,
        color=COLORS["warning"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(slide, "Screenshots captured 12 Aug 2026 from the production build · 559 unit tests, lint and build passing")


def slide_13(prs: Presentation):
    slide = new_slide(
        prs,
        14,
        "Value proposition",
        "The difference is what happens after the score",
        "FutureMe changes a result into a testable next step.",
    )
    # Traditional.
    add_pill(slide, "TRADITIONAL CAREER TEST", 0.78, 1.72, 2.35, fill=COLORS["surface2"], text_color=COLORS["muted"], line=COLORS["stroke"], size=9.5)
    traditional = [("Questionnaire", 2.42), ("Score / type", 3.42), ("Career list", 4.42), ("Stops", 5.42)]
    for label, y in traditional:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.98, y, 3.75, 0.62, fill=COLORS["surface"], line=COLORS["stroke"])
        add_text(slide, label, 1.18, y + 0.16, 3.35, 0.26, size=14, color=COLORS["muted"], bold=True, align=PP_ALIGN.CENTER)
        if y < 5.4:
            add_text(slide, "↓", 2.6, y + 0.63, 0.5, 0.34, size=18, color=COLORS["muted2"], bold=True, align=PP_ALIGN.CENTER)

    # FutureMe.
    add_pill(slide, "FUTUREME EVIDENCE LOOP", 5.35, 1.72, 2.35, fill=COLORS["mint"], text_color=COLORS["black"], size=9.5)
    future = [
        ("Reflect", "provisional signal", COLORS["indigo"]),
        ("Try", "independent evidence", COLORS["magenta"]),
        ("Explain", "reasons + unknowns", COLORS["mint"]),
        ("Compare", "constraints + trade-offs", COLORS["warning"]),
        ("Act", "30-day experiment", COLORS["coral"]),
    ]
    x_positions = [5.35, 6.85, 8.35, 9.85, 11.35]
    for i, ((name, detail, color), x) in enumerate(zip(future, x_positions)):
        add_circle_label(slide, str(i + 1), x + 0.36, 2.48, 0.58, fill=color, text_color=COLORS["black"] if color in (COLORS["mint"], COLORS["warning"]) else COLORS["white"], size=12)
        add_text(slide, name, x, 3.25, 1.28, 0.24, size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, detail, x - 0.05, 3.62, 1.38, 0.48, size=10.5, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        if i < 4:
            add_arrow_text(slide, x + 1.24, 2.59, COLORS["muted2"], 14)
    add_line(slide, 11.95, 4.22, 6.0, 4.22, color=COLORS["mint"], width=1.6, dash=True)
    add_text(slide, "new evidence updates the next exploration", 7.1, 4.28, 4.0, 0.24, size=10.5, color=COLORS["mint"], italic=True, align=PP_ALIGN.CENTER)

    differentiators = [
        "Refuses to guess",
        "Shows disagreement",
        "No manufactured winner",
        "Sources + freshness",
        "Vocational routes are first-class",
    ]
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.35, 5.08, 7.02, 1.18, fill=COLORS["surface2"], line=COLORS["mint"])
    dx = [5.55, 6.95, 8.35, 9.75, 11.15]
    for label, x in zip(differentiators, dx):
        add_text(slide, "✓", x, 5.3, 0.36, 0.26, size=15, color=COLORS["mint"], bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x - 0.35, 5.66, 1.12, 0.42, size=8.8, color=COLORS["text"], bold=True, align=PP_ALIGN.CENTER)
    add_source(slide, "Source: product design and implemented behavior")


def slide_14(prs: Presentation):
    slide = new_slide(
        prs,
        15,
        "Validation",
        "Validation is the next product—not a footnote",
        "The analysis pipeline is verified. The instrument is not.",
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 1.67, 4.05, 0.7, fill=COLORS["surface2"], line=COLORS["warning"])
    add_rich_line(
        slide,
        [
            ("PIPELINE VERIFIED", {"size": 15, "color": COLORS["mint"], "bold": True}),
            ("  ≠  ", {"size": 15, "color": COLORS["muted"], "bold": True}),
            ("INSTRUMENT VALIDATED", {"size": 15, "color": COLORS["warning"], "bold": True}),
        ],
        0.92,
        1.91,
        3.66,
        0.28,
        align=PP_ALIGN.CENTER,
    )
    steps = [
        ("0", "Ethics gate", "Approval · parent consent · student assent · PDPA", COLORS["coral"]),
        ("1", "Cognitive debrief", "n ≈ 30–40 · meaning and translation", COLORS["indigo"]),
        ("2", "Pilot", "target n = 250 · minimum 200", COLORS["magenta"]),
        ("3", "Measurement", "α + ω · item-total · “not sure” rate", COLORS["mint"]),
        ("4", "Structure", "RIASEC circular-order randomisation test", COLORS["warning"]),
        ("5", "Revise + re-pilot", "then invariance and outcome studies", COLORS["coral"]),
    ]
    y = 2.67
    for num, title, detail, color in steps:
        add_circle_label(slide, num, 0.85, y, 0.5, fill=color, text_color=COLORS["black"] if color in (COLORS["mint"], COLORS["warning"]) else COLORS["white"], size=11)
        add_text(slide, title, 1.55, y + 0.02, 2.2, 0.24, size=13.5, color=COLORS["text"], bold=True)
        add_text(slide, detail, 3.62, y + 0.02, 3.15, 0.28, size=10.7, color=COLORS["muted"])
        if num != "5":
            add_line(slide, 1.1, y + 0.5, 1.1, y + 0.72, color=COLORS["stroke"], width=1.3)
        y += 0.65

    add_section_label(slide, "PROPOSED SUCCESS MEASURES", 7.55, 1.78, 4.4, COLORS["mint"])
    measure_rows = [
        ("ASSESSMENT", "Reliability · item quality · construct structure"),
        ("ROUTES", "Relevance · counsellor agreement · diversity"),
        ("EXPERIENCE", "Completion · comprehension · route exploration"),
        ("LONG-TERM", "Can the student still explain and defend the choice?"),
    ]
    y = 2.28
    for cap, text_value in measure_rows:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.55, y, 4.75, 0.84, fill=COLORS["surface"], line=COLORS["stroke"])
        add_text(slide, cap, 7.82, y + 0.16, 1.05, 0.2, size=8.2, color=COLORS["mint"], bold=True, tracking=0.12)
        add_text(slide, text_value, 8.95, y + 0.15, 3.08, 0.48, size=11.5, color=COLORS["text"], bold=True)
        y += 1.02
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.55, 6.46, 4.75, 0.42, fill=COLORS["surface2"], line=COLORS["warning"])
    add_text(slide, "No reliability, norms or validity results exist yet.", 7.75, 6.57, 4.35, 0.2, size=10.5, color=COLORS["warning"], bold=True, align=PP_ALIGN.CENTER)
    add_source(slide, "Sources: FutureMe validation plan + pre-registered pilot protocol (30 Jul 2026)")


def slide_15(prs: Presentation):
    slide = new_slide(
        prs,
        16,
        "Roadmap",
        "From runnable demo to trusted decision support",
        "Every phase has a gate; partnerships stay exploratory until something is signed.",
    )
    phases = [
        ("1", "Foundation", "DONE", COLORS["mint"]),
        ("2", "Runnable\nprototype", "DONE", COLORS["mint"]),
        ("3", "Validation", "NOW · NEXT", COLORS["warning"]),
        ("4", "School pilot", "PLANNED", COLORS["magenta"]),
        ("5", "Ecosystem", "EXPLORATORY", COLORS["indigo"]),
    ]
    x_positions = [0.8, 3.1, 5.4, 7.7, 10.0]
    add_line(slide, 1.35, 2.2, 10.55, 2.2, color=COLORS["stroke"], width=2.0)
    for num, name, status, color in phases:
        x = x_positions[int(num) - 1]
        add_circle_label(slide, num, x, 1.73, 0.94, fill=color, text_color=COLORS["black"] if color in (COLORS["mint"], COLORS["warning"]) else COLORS["white"], size=16)
        add_text(slide, name, x - 0.4, 2.85, 1.75, 0.52, size=14, color=COLORS["text"], bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, status, x - 0.5, 3.45, 1.95, 0.2, size=8.2, color=color, bold=True, tracking=0.22, align=PP_ALIGN.CENTER)

    add_text(
        slide,
        "Phase 3: adapt Thai translation · validate instrument + missions · licensed route data · bias + safety audit",
        1.05,
        4.02,
        11.2,
        0.42,
        size=12.5,
        color=COLORS["warning"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.78, 4.65, 5.72, 0.78, fill=COLORS["surface2"], line=COLORS["mint"])
    add_rich_line(
        slide,
        [
            ("Career test  →  ", {"size": 15.5, "color": COLORS["muted"]}),
            ("AI career discovery companion", {"size": 15.5, "color": COLORS["mint"], "bold": True}),
        ],
        1.05,
        4.9,
        5.2,
        0.3,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "ช่วยให้เด็กเข้าใจตัวเอง ทดลองทางเลือก และตัดสินใจด้วยเหตุผลที่อธิบายได้",
        0.98,
        5.56,
        5.35,
        0.42,
        size=14,
        color=COLORS["text"],
        bold=True,
        font=FONT_THAI,
        align=PP_ALIGN.CENTER,
    )

    add_section_label(slide, "SELECTED REFERENCES", 6.9, 4.64, 4.0, COLORS["mint"])
    refs = [
        "[1] TDRI — การพัฒนาทุนมนุษย์ไทย (2025)",
        "[2] NSO — Social Indicators 2025, p.185",
        "[3] NESDC — ภาวะสังคมไทย Q2/2568",
        "[4] OECD — PIAAC 2023 mismatch evidence",
        "[5] WEF — Future of Jobs 2025",
        "[6] Holland (1997) + O*NET Interest Profiler Manual",
        "[7] FutureMe Source Audit, Methodology & Pilot Protocol",
    ]
    add_paragraphs(slide, refs, 6.9, 5.02, 5.55, 1.68, size=9.5, color=COLORS["muted"], spacing=2.5)
    add_source(slide, "Full source registry: 01_Research/Data/REFERENCES.md")


def build_deck() -> Path:
    make_backgrounds()
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "FutureMe AI — Project Case Study"
    prs.core_properties.subject = "Problem, evidence, design, implementation, validation and roadmap"
    prs.core_properties.author = "FutureMe AI project team"
    prs.core_properties.keywords = "FutureMe, Thai education, career exploration, RIASEC, explainable AI"
    prs.core_properties.comments = (
        "Generated from the audited Data folder and the runnable prototype. "
        "Current versus planned capabilities are separated throughout."
    )

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)
    slide_11(prs)
    slide_11b(prs)
    slide_12(prs)
    slide_13(prs)
    slide_14(prs)
    slide_15(prs)

    prs.save(PPTX_PATH)
    return PPTX_PATH


if __name__ == "__main__":
    result = build_deck()
    print(result)
